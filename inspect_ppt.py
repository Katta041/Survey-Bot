from pptx import Presentation
import os

file_path = '/Users/aierarohit/Desktop/Political Data/Viralimalai_Baseline Survey Report_V1.pptx'

print(f"Inspecting file: {file_path}")

try:
    prs = Presentation(file_path)
    print(f"Number of slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        print("\n".join(text_content))

except Exception as e:
    print(f"Error inspecting PPT file: {e}")
